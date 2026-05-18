"""
Genera la presentazione PowerPoint da mostrare al Direttivo per
illustrare il nuovo sito orbassanocalcio.com.

Output: ./presentazione-direttivo.pptx (gitignored).

Esecuzione: python scripts/generate-deck-direttivo.py
Richiede: python-pptx (gia' installato nel sistema, vedi
extract-palette.py).

Brand:
  navy   #0a1428 (sfondo principale)
  blue   #213f8c (accent secondario)
  red    #e91f22 (accent CTA)
  gold   #dfb16c (highlight / eyebrow)
  ivory  #fefdfd (testo principale su navy)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- Palette --------------------------------------------------------------
NAVY = RGBColor(0x0A, 0x14, 0x28)
BLUE = RGBColor(0x21, 0x3F, 0x8C)
RED = RGBColor(0xE9, 0x1F, 0x22)
GOLD = RGBColor(0xDF, 0xB1, 0x6C)
IVORY = RGBColor(0xFE, 0xFD, 0xFD)
INK_MID = RGBColor(0xA8, 0xB5, 0xCC)
INK_LOW = RGBColor(0x6B, 0x7A, 0x99)

# 16:9 standard (13.333 x 7.5 inch)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_bg(slide, color=NAVY):
    """Sfondo a tinta unita per la slide."""
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    rect.shadow.inherit = False
    return rect


def add_text(
    slide,
    text,
    left,
    top,
    width,
    height,
    *,
    font="Calibri",
    size=Pt(18),
    color=IVORY,
    bold=False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    tracking=0,
):
    """Helper per aggiungere un blocco di testo."""
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = size
    r.font.bold = bold
    r.font.color.rgb = color
    return tx, tf, p, r


def add_eyebrow(slide, text, left, top, color=GOLD):
    """Eyebrow piccolo, gold, uppercase tracking ampio."""
    add_text(
        slide,
        text.upper(),
        left,
        top,
        Inches(8),
        Inches(0.35),
        font="Calibri",
        size=Pt(12),
        color=color,
        bold=True,
        align=PP_ALIGN.LEFT,
    )


def add_h1(slide, text, left, top, width=None):
    """Titolo grande slide."""
    if width is None:
        width = SLIDE_W - left - Inches(0.8)
    return add_text(
        slide,
        text,
        left,
        top,
        width,
        Inches(1.6),
        font="Calibri",
        size=Pt(44),
        color=IVORY,
        bold=True,
        align=PP_ALIGN.LEFT,
    )


def add_underline(slide, left, top, width=Inches(0.8), color=GOLD):
    """Linea oro orizzontale 4px, accento sotto eyebrow/titoli."""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Emu(38100)  # ~3pt
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def add_footer(slide, slide_num, total):
    """Footer minimal: nome club a sinistra, numero pagina a destra."""
    add_text(
        slide,
        "ASD ORBASSANO CALCIO  ·  Sito ufficiale 2026",
        Inches(0.6),
        Inches(7.0),
        Inches(7),
        Inches(0.4),
        font="Calibri",
        size=Pt(9),
        color=INK_LOW,
        bold=False,
        tracking=200,
    )
    add_text(
        slide,
        f"{slide_num:02d} / {total:02d}",
        SLIDE_W - Inches(1.6),
        Inches(7.0),
        Inches(1),
        Inches(0.4),
        font="Calibri",
        size=Pt(9),
        color=GOLD,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )


# --- Build deck -----------------------------------------------------------
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # blank

TOTAL = 8


# --- SLIDE 1: COVER -------------------------------------------------------
s = prs.slides.add_slide(blank_layout)
add_bg(s, NAVY)

# Striscia diagonale gold sullo sfondo (accent visivo) — top-right corner
accent = s.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(11),
    Inches(-1),
    Inches(0.18),
    Inches(4.5),
)
accent.fill.solid()
accent.fill.fore_color.rgb = GOLD
accent.line.fill.background()
accent.rotation = 25

add_eyebrow(s, "Presentazione al Direttivo  ·  Maggio 2026", Inches(0.8), Inches(0.8))
add_underline(s, Inches(0.8), Inches(1.18))

add_text(
    s,
    "Il nuovo sito ufficiale",
    Inches(0.8),
    Inches(2.4),
    Inches(11),
    Inches(0.8),
    font="Calibri",
    size=Pt(28),
    color=INK_MID,
    bold=False,
)
add_text(
    s,
    "ORBASSANOCALCIO.COM",
    Inches(0.8),
    Inches(3.1),
    Inches(12),
    Inches(2),
    font="Calibri",
    size=Pt(80),
    color=IVORY,
    bold=True,
)
add_text(
    s,
    "Dal 1930 il rossoblù di Orbassano  ·  Dal 2026 anche online.",
    Inches(0.8),
    Inches(5.2),
    Inches(11),
    Inches(0.6),
    font="Calibri",
    size=Pt(18),
    color=GOLD,
    bold=False,
)

add_text(
    s,
    "Costruito su Next.js 16 + Sanity CMS  ·  Hosting Vercel",
    Inches(0.8),
    Inches(6.4),
    Inches(11),
    Inches(0.4),
    font="Calibri",
    size=Pt(11),
    color=INK_LOW,
)


# --- SLIDE 2: COSA ABBIAMO COSTRUITO --------------------------------------
s = prs.slides.add_slide(blank_layout)
add_bg(s, NAVY)
add_eyebrow(s, "01  ·  Il progetto", Inches(0.8), Inches(0.6))
add_underline(s, Inches(0.8), Inches(0.95))
add_h1(s, "Cosa abbiamo costruito", Inches(0.8), Inches(1.3))

bullets = [
    ("Sito completo e moderno", "32 pagine pubbliche, design custom, mobile-first."),
    ("CMS gestito dal club", "Il direttivo aggiorna tutto in autonomia, senza dipendere da agenzie."),
    ("Zero lock-in", "Codice sorgente, dati e dominio appartengono al club. Nessuna piattaforma chiusa."),
    ("Performance enterprise-grade", "Hosting CDN globale, caricamento sotto 1 secondo, sicurezza HTTPS."),
]
top = Inches(2.9)
for i, (head, body) in enumerate(bullets):
    add_text(
        s, head, Inches(0.8), top, Inches(11.5), Inches(0.45),
        font="Calibri", size=Pt(20), color=IVORY, bold=True,
    )
    add_text(
        s, body, Inches(0.8), top + Inches(0.45), Inches(11.5), Inches(0.45),
        font="Calibri", size=Pt(14), color=INK_MID,
    )
    top += Inches(1.0)

add_footer(s, 2, TOTAL)


# --- SLIDE 3: SEZIONI DEL SITO --------------------------------------------
s = prs.slides.add_slide(blank_layout)
add_bg(s, NAVY)
add_eyebrow(s, "02  ·  Mappa del sito", Inches(0.8), Inches(0.6))
add_underline(s, Inches(0.8), Inches(0.95))
add_h1(s, "Le sezioni live", Inches(0.8), Inches(1.3))

# 3 colonne di sezioni
cols = [
    ("SQUADRE & CALENDARI", [
        "Prima Squadra · Promozione",
        "Juniores Under 19",
        "Settore Giovanile (U17 → U14)",
        "Calendario aggregato SGS",
        "Archivio stagioni passate",
    ]),
    ("EDITORIALE", [
        "News & comunicati",
        "Gallery foto e album",
        "Storia del club (timeline 1930)",
        "Organigramma e staff",
        "Impianti sportivi",
    ]),
    ("OPERATIVO", [
        "Open Days e Tornei",
        "Modulo iscrizione SGS",
        "Sponsor & Partner",
        "Biglietteria",
        "5×1000, Newsletter, Contatti",
    ]),
]

col_w = Inches(4.0)
col_x = [Inches(0.8), Inches(4.85), Inches(8.9)]
for i, (title, items) in enumerate(cols):
    add_text(
        s, title, col_x[i], Inches(2.9), col_w, Inches(0.4),
        font="Calibri", size=Pt(12), color=GOLD, bold=True,
    )
    add_underline(s, col_x[i], Inches(3.25), Inches(0.5))
    top = Inches(3.55)
    for item in items:
        add_text(
            s, "·  " + item, col_x[i], top, col_w, Inches(0.4),
            font="Calibri", size=Pt(14), color=IVORY,
        )
        top += Inches(0.5)

add_footer(s, 3, TOTAL)


# --- SLIDE 4: TOOLS PER IL CLUB -------------------------------------------
s = prs.slides.add_slide(blank_layout)
add_bg(s, NAVY)
add_eyebrow(s, "03  ·  Strumenti operativi", Inches(0.8), Inches(0.6))
add_underline(s, Inches(0.8), Inches(0.95))
add_h1(s, "Cosa può fare il club da solo", Inches(0.8), Inches(1.3))

tools = [
    ("CMS Sanity Studio", "Interfaccia in italiano per modificare news, foto, rose, calendari, sponsor. Modifiche live in 60 secondi."),
    ("Import massivo da Excel", "Una sola operazione carica 100+ giocatori o un intero calendario di stagione. Riusabile ogni anno."),
    ("Newsletter integrata", "Form pubblico con conferma email. Lista gestita su Brevo, primo invio in 5 minuti."),
    ("Email transazionali", "Contatti, segnalazioni, iscrizioni SGS recapitate via Resend, niente caselle PEC piene."),
    ("Webhook automatici", "Ogni modifica in CMS rigenera la pagina sul sito in pochi secondi, senza redeploy manuale."),
]
top = Inches(2.9)
for head, body in tools:
    add_text(
        s, head, Inches(0.8), top, Inches(11.5), Inches(0.4),
        font="Calibri", size=Pt(17), color=GOLD, bold=True,
    )
    add_text(
        s, body, Inches(0.8), top + Inches(0.4), Inches(11.5), Inches(0.4),
        font="Calibri", size=Pt(13), color=INK_MID,
    )
    top += Inches(0.82)

add_footer(s, 4, TOTAL)


# --- SLIDE 5: NUMERI ------------------------------------------------------
s = prs.slides.add_slide(blank_layout)
add_bg(s, NAVY)
add_eyebrow(s, "04  ·  Il sito in cifre", Inches(0.8), Inches(0.6))
add_underline(s, Inches(0.8), Inches(0.95))
add_h1(s, "Numeri al lancio", Inches(0.8), Inches(1.3))

stats = [
    ("107", "giocatori in rosa"),
    ("6", "squadre attive"),
    ("26+", "partite caricate"),
    ("32", "pagine pubbliche"),
    ("13", "club avversari"),
    ("< 1s", "tempo medio di caricamento"),
]
# Griglia 3x2
cell_w = Inches(3.9)
cell_h = Inches(1.6)
gap = Inches(0.25)
start_x = Inches(0.8)
start_y = Inches(3.0)
for idx, (num, label) in enumerate(stats):
    col = idx % 3
    row = idx // 3
    x = start_x + (cell_w + gap) * col
    y = start_y + (cell_h + gap) * row
    # cell border-left gold
    rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Emu(38100), cell_h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = GOLD
    rect.line.fill.background()
    # number
    add_text(
        s, num, x + Inches(0.25), y - Inches(0.05), cell_w, Inches(1.0),
        font="Calibri", size=Pt(56), color=IVORY, bold=True,
    )
    # label
    add_text(
        s, label.upper(), x + Inches(0.25), y + Inches(1.05), cell_w, Inches(0.4),
        font="Calibri", size=Pt(11), color=INK_MID, bold=True,
    )

add_footer(s, 5, TOTAL)


# --- SLIDE 6: COSTI ANNUALI -----------------------------------------------
s = prs.slides.add_slide(blank_layout)
add_bg(s, NAVY)
add_eyebrow(s, "05  ·  Budget di gestione", Inches(0.8), Inches(0.6))
add_underline(s, Inches(0.8), Inches(0.95))
add_h1(s, "Costi annuali ricorrenti", Inches(0.8), Inches(1.3))

costs = [
    ("Hosting Vercel", "Free → Pro", "€ 0 – 240 / anno"),
    ("CMS Sanity", "Free → Growth", "€ 0 – 300 / anno"),
    ("Email transazionali (Resend)", "Free tier 3.000 / mese", "€ 0 – 120 / anno"),
    ("Newsletter (Brevo)", "Free tier 300 / giorno", "€ 0 – 60 / anno"),
    ("Dominio orbassanocalcio.com", "registrar standard", "≈ € 15 / anno"),
]
top = Inches(2.9)
for service, plan, price in costs:
    # service
    add_text(
        s, service, Inches(0.8), top, Inches(5.5), Inches(0.4),
        font="Calibri", size=Pt(15), color=IVORY, bold=True,
    )
    # plan
    add_text(
        s, plan, Inches(6.4), top, Inches(4), Inches(0.4),
        font="Calibri", size=Pt(13), color=INK_MID,
    )
    # price
    add_text(
        s, price, Inches(10.5), top, Inches(2.3), Inches(0.4),
        font="Calibri", size=Pt(15), color=GOLD, bold=True, align=PP_ALIGN.RIGHT,
    )
    # divider
    div = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), top + Inches(0.5), Inches(12), Emu(6350)
    )
    div.fill.solid()
    div.fill.fore_color.rgb = INK_LOW
    div.line.fill.background()
    top += Inches(0.65)

# Total
total_top = top + Inches(0.3)
add_text(
    s, "TOTALE", Inches(0.8), total_top, Inches(5.5), Inches(0.5),
    font="Calibri", size=Pt(16), color=IVORY, bold=True,
)
add_text(
    s, "€ 15 – 735 / anno", Inches(10.5), total_top, Inches(2.3), Inches(0.5),
    font="Calibri", size=Pt(20), color=GOLD, bold=True, align=PP_ALIGN.RIGHT,
)

add_footer(s, 6, TOTAL)


# --- SLIDE 7: VALUTAZIONE DI MERCATO --------------------------------------
s = prs.slides.add_slide(blank_layout)
add_bg(s, NAVY)
add_eyebrow(s, "06  ·  Valutazione patrimoniale", Inches(0.8), Inches(0.6))
add_underline(s, Inches(0.8), Inches(0.95))
add_h1(s, "Quanto vale, sul mercato", Inches(0.8), Inches(1.3))

add_text(
    s,
    "Riferimenti agenzie / studi del settore in Italia, anno 2026:",
    Inches(0.8), Inches(2.9), Inches(11), Inches(0.4),
    font="Calibri", size=Pt(14), color=INK_MID,
)

tiers = [
    ("Sito club Serie D / dilettanti", "WordPress + tema", "€ 3.000 – 8.000"),
    ("Equivalente agenzia digital media-tier", "Next.js + headless CMS, no design custom", "€ 15.000 – 35.000"),
    ("Equivalente boutique specializzata club", "design + brand system + UX research", "€ 45.000 – 80.000"),
    ("Sito club Serie A / B", "CRM, ticketing, e-commerce, multilingua", "€ 100.000 +"),
]
top = Inches(3.5)
for tier, detail, price in tiers:
    add_text(
        s, tier, Inches(0.8), top, Inches(5.5), Inches(0.4),
        font="Calibri", size=Pt(14), color=IVORY, bold=True,
    )
    add_text(
        s, detail, Inches(0.8), top + Inches(0.4), Inches(5.5), Inches(0.35),
        font="Calibri", size=Pt(11), color=INK_LOW,
    )
    add_text(
        s, price, Inches(8.5), top + Inches(0.15), Inches(4.3), Inches(0.45),
        font="Calibri", size=Pt(18), color=GOLD, bold=True, align=PP_ALIGN.RIGHT,
    )
    top += Inches(0.85)

# Box riassuntivo
box = s.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(6.85), Inches(12), Inches(0.6)
)
box.fill.solid()
box.fill.fore_color.rgb = BLUE
box.line.fill.background()
add_text(
    s,
    "Stima realistica per il nostro sito chiavi in mano:  € 20.000 – 35.000  ·  asset livello Serie C",
    Inches(0.9), Inches(6.92), Inches(12), Inches(0.5),
    font="Calibri", size=Pt(13), color=IVORY, bold=True,
)


# --- SLIDE 8: PROSSIMI PASSI ----------------------------------------------
s = prs.slides.add_slide(blank_layout)
add_bg(s, NAVY)
add_eyebrow(s, "07  ·  Roadmap", Inches(0.8), Inches(0.6))
add_underline(s, Inches(0.8), Inches(0.95))
add_h1(s, "Prossimi passi & supporto del Direttivo", Inches(0.8), Inches(1.3))

# Colonna sinistra: cosa serve dal direttivo
add_text(
    s, "COSA SERVE DAL DIRETTIVO", Inches(0.8), Inches(2.9), Inches(5.5), Inches(0.4),
    font="Calibri", size=Pt(12), color=GOLD, bold=True,
)
add_underline(s, Inches(0.8), Inches(3.25))
left_items = [
    "Approvazione testi pubblici (Codice Etico, Trasparenza)",
    "Caricamento foto giocatori e team (server: Cloudinary)",
    "Conferma migrazione DNS da Wix",
    "Validazione contenuti storia, organigramma, impianti",
    "Approvazione comunicazione di lancio + 5×1000",
]
top = Inches(3.55)
for item in left_items:
    add_text(
        s, "·  " + item, Inches(0.8), top, Inches(5.7), Inches(0.5),
        font="Calibri", size=Pt(13), color=IVORY,
    )
    top += Inches(0.55)

# Colonna destra: prossimi step tecnici
add_text(
    s, "PROSSIMI STEP TECNICI", Inches(7.0), Inches(2.9), Inches(5.5), Inches(0.4),
    font="Calibri", size=Pt(12), color=GOLD, bold=True,
)
add_underline(s, Inches(7.0), Inches(3.25))
right_items = [
    "Setup analytics privacy-friendly (no Google)",
    "Switch DNS orbassanocalcio.com → Vercel",
    "Smoke test post-DNS (form, newsletter, SEO)",
    "Lancio ufficiale + comunicato sui social",
    "Tracciamento Search Console + indicizzazione",
]
top = Inches(3.55)
for item in right_items:
    add_text(
        s, "·  " + item, Inches(7.0), top, Inches(5.7), Inches(0.5),
        font="Calibri", size=Pt(13), color=IVORY,
    )
    top += Inches(0.55)

# CTA finale
cta = s.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.6), Inches(12), Inches(0.7)
)
cta.fill.solid()
cta.fill.fore_color.rgb = RED
cta.line.fill.background()
add_text(
    s,
    "Pronti al lancio. Restano da chiudere DNS + ultimi contenuti CMS.",
    Inches(0.8), Inches(6.7), Inches(12), Inches(0.5),
    font="Calibri", size=Pt(15), color=IVORY, bold=True, align=PP_ALIGN.CENTER,
)


# --- Save -----------------------------------------------------------------
output_path = "presentazione-direttivo.pptx"
prs.save(output_path)

import os
size_kb = os.path.getsize(output_path) / 1024
print(f"\nOK - Generata: {output_path}")
print(f"  Dimensione: {size_kb:.1f} KB")
print(f"  Slide: {len(prs.slides)}")
print(f"\nApri con PowerPoint, Keynote, Google Slides o LibreOffice Impress.\n")
